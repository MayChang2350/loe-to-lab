import glob, os
from pptx import Presentation
outdir="/sessions/relaxed-nice-hopper/mnt/outputs/slidetext"
os.makedirs(outdir, exist_ok=True)
for f in sorted(glob.glob("/sessions/relaxed-nice-hopper/mnt/uploads/*.pptx")):
    base=os.path.basename(f)
    try:
        prs=Presentation(f)
    except Exception as e:
        print("ERR",base,e); continue
    lines=[f"=== {base} ({len(prs.slides)} slides) ==="]
    for i,s in enumerate(prs.slides,1):
        lines.append(f"--- slide {i} ---")
        for sh in s.shapes:
            if sh.has_text_frame:
                t=sh.text_frame.text.strip()
                if t: lines.append(t)
            if sh.has_table:
                for r in sh.table.rows:
                    lines.append(" | ".join(c.text.strip() for c in r.cells))
        try:
            if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
                lines.append("[NOTES] "+s.notes_slide.notes_slide if False else "[NOTES] "+s.notes_slide.notes_text_frame.text.strip())
        except Exception: pass
    open(os.path.join(outdir, base+".txt"),"w").write("\n".join(lines))
    print(base, len(prs.slides))
